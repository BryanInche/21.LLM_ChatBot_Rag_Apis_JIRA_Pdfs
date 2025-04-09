"""
1. MÓDULO DE CARGA Y PROCESAMIENTO DEL VECTOR STORE DE LOS DOCUMENTOS RAG 
=====================================================

Este módulo proporciona funcionalidades para:
1. Cargar documentos en varios formatos (PDF, PPT, DOCX)
2. Dividir los documentos en fragmentos(chunks) adecuados para procesamiento
3. Crear y gestionar un vector store (ChromaDB) con embeddings de los documentos

Dependencias clave:
- LangChain: Framework para aplicaciones con modelos de lenguaje
- PyPDF/PyPDF2: Para procesamiento de PDFs
- python-pptx: Para procesamiento de PowerPoint
- python-docx: Para procesamiento de Word
- ChromaDB: Base de datos vectorial
- HuggingFace Embeddings: Modelos de embeddings
"""
import os
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings.fastembed import FastEmbedEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_huggingface import HuggingFaceEmbeddings
#pip install python-pptx
#pip install python-docx  
#pip install  
from pptx import Presentation
from docx import Document
#from langchain.docstore.document import Document as LangchainDocument  # Importar el objeto Document de LangChain
import pypandoc
#from atlassian import Jira
#pip install langchain-unstructured

from langchain_core.documents import Document as LangchainDocument  # Versión moderna
import json
import requests  # hacer llamadas HTTP, HTTPBasicAuth para autenticación
from requests.auth import HTTPBasicAuth 


def obtener_datos_jira(jql, campos, max_results=100):
    """Obtiene datos de JIRA API"""
    # Configura tus datos
    JIRA_DOMAIN = "ms4m.atlassian.net"
    EMAIL = "bryan.inche@ms4m.com"
    API_TOKEN = "ATATT3xFfGF0s-RwXi3ojpNz7ZzommvXwJh18Xq1KftfDGNobAJyT8URut8evud2l1JcSMK2xUK9A4KzNZoXMrLlmtg6XGEVQt2V8ErdHVUrk4UYvcoHC6W4oqfaQ--48_dmITFcDZZiTqciN1XqtykkHYP64AYyS-viZyzHnBEPDs07KNC7wlc=A7B98677"

    # Endpoint de Jira
    url = f"https://{JIRA_DOMAIN}/rest/api/2/search"
    all_issues = []
    start_at = 0
    auth = HTTPBasicAuth(EMAIL, API_TOKEN)

    while True:
        response = requests.post(
            url=url,
            headers={"Content-Type": "application/json"},
            auth=auth,
            json={
                "jql": jql,
                "startAt": start_at,
                "maxResults": max_results,
                "fields": campos  # Los campos se pasan directamente
            }
        )
        
        #response.raise_for_status()
        data = response.json()
        issues = data.get("issues", [])
        all_issues.extend(issues)
        
        if len(issues) < max_results:
            break
        start_at += max_results
        print(f"Descargados {len(issues)} tickets (total: {len(all_issues)})")
    
    return all_issues



def jira_issues_a_documentos(issues):
    """
    Convierte issues de JIRA(JSON) a documentos LangChain
    (Versión compatible con tu estructura actual)
    
    Args:
        issues (list): Issues de JIRA en formato JSON
        
    Returns:
        list: Lista de objetos Document de LangChain con:
              - page_content: Texto estructurado del ticket
              - metadata: Metadatos para filtrado/búsqueda
    """
    documentos = []

    for ticket in issues:
        # Extracción de campos principales
        key = ticket.get("key", "SIN CLAVE")
        fields_data = ticket.get("fields", {})
        
        summary = fields_data.get("summary", "")
        description = fields_data.get("description", "")
        status = fields_data.get("status", {}).get("name", "")
        resolution = fields_data.get("resolution", {}).get("name", "")
        resolutiondate = fields_data.get("resolutiondate", "")
        solucion_final = fields_data.get("customfield_10975", "")
        causa_raiz = fields_data.get("customfield_11282", "")
        
        # Campos personalizados adicionales (como en tu código)
        causa_immediata = fields_data.get("customfield_11281", "")
        #customfield_11228 = fields_data.get("customfield_11228", "")
        

        # Construir contenido principal (similar a tu formato)
        contenido = f"""
        Resumen: {summary}
        Descripción: {description}
        Estado: {status}
        Resolución: {resolution}
        Fecha de resolución: {resolutiondate}
        Solución final: {solucion_final}
        Causa Raiz: {causa_raiz}
        Causa Immediata: {causa_immediata}
        """.strip()

        # Documento principal (misma estructura de metadata)
        documentos.append(LangchainDocument(
            page_content=contenido,
            metadata={
                "ticket": key,
                "tipo": "contexto_completo",
                "status": status,
                "resolution": resolution,
                "resolution_date": resolutiondate
            }
        ))

        # Procesar comentarios (igual a tu implementación)
        comentarios = fields_data.get("comment", {}).get("comments", [])
        for idx, comentario in enumerate(comentarios):
            body = comentario.get("body", "")
            if body:
                documentos.append(LangchainDocument(
                    page_content=body.strip(),
                    metadata={
                        "ticket": key,
                        "tipo": f"comentario_{idx+1}",
                        "autor": comentario.get("author", {}).get("displayName", "Anónimo"),
                        "fecha": comentario.get("created", "")
                    }
                ))
    
    return documentos



def cargar_documentos(rutas_archivos, jira_config=None):
    """
    Carga y procesa documentos desde múltiples fuentes y formatos.
    
    Funcionamiento detallado:
    1. Acepta rutas individuales, listas de rutas o directorios
    2. Detecta automáticamente el tipo de archivo (PDF, PPT, DOCX)
    3. Extrae texto usando el parser apropiado para cada formato
    4. Divide el contenido en fragmentos manejables
    
    Parámetros:
        rutas_archivos (str/list): 
            - Ruta a un archivo individual (str)
            - Lista de rutas a archivos (list)
            - Ruta a un directorio (se procesarán todos los archivos compatibles)
    
    Retorna:
        list: Lista de documentos divididos en fragmentos.
    """
    todos_los_documentos = []

    # Si se proporciona una sola ruta (str), convertirla en una lista
    if isinstance(rutas_archivos, str):
        rutas_archivos = [rutas_archivos]

    for ruta in rutas_archivos:
        # Procesar directorios: Si es una carpeta, obtener todos los archivos compatibles dentro de ella
        if os.path.isdir(ruta):
            print(f"Procesando carpeta: {ruta}")
            archivos_en_carpeta = [
                os.path.join(ruta, archivo) for archivo in os.listdir(ruta)
                if archivo.endswith(".pdf") or archivo.endswith(".ppt") or archivo.endswith(".docx") or archivo.endswith(".json")
            ]
            if not archivos_en_carpeta:
                print(f"Advertencia: No se encontraron archivos compatibles en la carpeta {ruta}.")
            else:
                print(f"Se encontraron {len(archivos_en_carpeta)} archivos compatibles en la carpeta {ruta}.")
            rutas_archivos.extend(archivos_en_carpeta)
            continue
        
        # Si es un archivo PDF, cargarlo
        if ruta.endswith(".pdf"):
            try:
                print(f"Procesando archivo PDF: {ruta}")
                carga = PyPDFLoader(ruta)
                #carga = UnstructuredFileLoader(ruta) # Para PDFs problemáticos
                documentos = carga.load()
            except Exception as e:
                print(f"Error al cargar el archivo PDF {ruta}: {e}")
                continue

        # Si es un archivo PPTX, cargarlo
        elif ruta.endswith(".ppt"):
            try:
                print(f"Procesando archivo PPT: {ruta}")
                ############################################
                #prs = Presentation(ruta)
                #textos = []
                #for slide in prs.slides:
                #    for shape in slide.shapes:
                #        if hasattr(shape, "text"):
                #            textos.append(shape.text)
                ##############################################
                texto = pypandoc.convert_file(ruta, 'plain') #extraer el texto directamente desde los archivos .ppt

                # Crear un objeto Document de LangChain
                documentos = [LangchainDocument(page_content=texto)]            
                # Crear un objeto Document de LangChain
                #documentos = [LangchainDocument(page_content="\n".join(textos))]
            except Exception as e:
                print(f"Error al cargar el archivo PPT {ruta}: {e}")
                continue

        # Si es un archivo DOCX, cargarlo
        elif ruta.endswith(".docx"):
            try:
                print(f"Procesando archivo DOCX: {ruta}")
                doc = Document(ruta)
                textos = [para.text for para in doc.paragraphs]
                # Crear un objeto Document de LangChain
                documentos = [LangchainDocument(page_content="\n".join(textos))]
            except Exception as e:
                print(f"Error al cargar el archivo DOCX {ruta}: {e}")
                continue
        
        # Si es un archivo JSON, cargarlo
        elif ruta.lower().endswith(".json"):
            try:
                print(f"Procesando archivo JIRA JSON: {ruta}")
                with open(ruta, "r", encoding="utf-8") as f:
                    issues = json.load(f)
                print(f"Tickets JIRA encontrados: {len(issues)}")
                #print("Ejemplo de ticket cargado:")
                #print(json.dumps(issues[0], indent=4, ensure_ascii=False))  # Muestra el primer ticket con formato legible

                documentos_jira = jira_issues_a_documentos(issues)

                print("Ejemplo de ticket cargado preprocesado:")
                if documentos_jira:
                    primer_doc = documentos_jira[0]
                    print(f"Contenido: {primer_doc.page_content[:200]}...")  # Muestra los primeros 200 caracteres
                    print("Metadatos:")
                    for k, v in primer_doc.metadata.items():
                        print(f"  {k}: {v}")
                    print("-" * 50)

                todos_los_documentos.extend(documentos_jira)
                continue  # Saltar el split para JIRA (ya viene estructurado)
            except Exception as e:
                print(f"Error al procesar JSON de JIRA: {e}")
                continue

        else:
            print(f"Advertencia: La ruta {ruta} no es un archivo compatible.")
            continue

        if not documentos:
            print(f"Advertencia: No se pudieron cargar documentos desde el archivo {ruta}.")
            continue

        # Dividir los documentos en fragmentos (excepto JIRA) (Json ya lo divide)
        try:
            texto_spliter = RecursiveCharacterTextSplitter(chunk_size=5000, chunk_overlap=800)
            documentos_divididos = texto_spliter.split_documents(documentos)
        except Exception as e:
            print(f"Error al dividir los documentos del archivo {ruta}: {e}")
            continue

        # Agregar los documentos divididos a la lista general
        todos_los_documentos.extend(documentos_divididos)
        print(f"Se cargaron {len(documentos_divididos)} fragmentos desde {ruta}.")

    # ===== 2. Carga desde API de JIRA (si se proporciona configuración) =====
    if jira_config:
        print("\n Conectando a API de JIRA...")
        try:
            issues = obtener_datos_jira(
                jql=jira_config['jql'],
                campos=jira_config['fields']
                #campos=jira_config.get('fields', ["key", "summary", "status", "resolution", "resolutiondate", "description", "comment", "customfield_11281",
                #"customfield_11228", "customfield_10975"])
            )
            print(f"Tickets obtenidos: {len(issues)}")
            
            documentos_jira = jira_issues_a_documentos(issues)
            print(f"Documentos generados: {len(documentos_jira)}")
            
            # Mostrar ejemplo (opcional)
            if documentos_jira:
                print("\n  [Ejemplo de ticket desde API]")
                primer_doc = documentos_jira[0]
                print(f"  Contenido (inicio): {primer_doc.page_content[:100]}...")
                print("  Metadatos:")
                for k, v in primer_doc.metadata.items():
                    print(f"    {k}: {v}")
            
            todos_los_documentos.extend(documentos_jira)
        except Exception as e:
            print(f"Error al obtener datos de JIRA: {str(e)}")


    if not todos_los_documentos:
        raise ValueError("No se pudieron cargar documentos desde ninguno de los archivos o carpetas proporcionados.")

    print(f"Se cargaron un total de {len(todos_los_documentos)} fragmentos de texto.")
    return todos_los_documentos



def crear_vectorstore(docs):
    """
    Crea y configura una base de datos vectorial Chroma con embeddings.
    
    Proceso detallado:
    1. Inicializa el modelo de embeddings (HuggingFace)
    2. Crea la base de datos Chroma con persistencia
    3. Valida la correcta creación del índice
    4. Realiza una consulta de prueba
    
    Parámetros:
        docs (list): Lista de documentos LangChain a indexar
    
    Retorna:
        Chroma: Instancia configurada del vector store creado. 
    """
    # Validar que se proporcionen documentos
    if not docs:
        raise ValueError("No se proporcionaron documentos para crear el vector store.")

    # Crear el modelo de embeddings
    try:
        # Busqueda semantica, basada en similitud de embedings
        # Modelo compatible con FastEmbedEmbeddings
        #embeding_modelo = FastEmbedEmbeddings(model_name="BAAI/bge-small-en-v1.5")
        
        # Modelo de Hugging Face
        embeding_modelo = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    except Exception as e:
        print(f"Error al crear el modelo de embeddings: {e}")
        raise

    # Crear el vector store
    try:
        vector_store = Chroma.from_documents(
            documents=docs,
            embedding=embeding_modelo,
            persist_directory="5chroma_ms4m_bd_pdf_jira",
            collection_name="5data_ms4m_pdf_jira" #colection dentro del vs
        )
        print("Vector store creado correctamente.")
    except Exception as e:
        print(f"Error al crear el vector store: {e}")
        raise

    # Depuración: Verificar el número de documentos en el vector store
    print(f"Número de documentos en el vector store: {vector_store._collection.count()}")

    # Validar que el vector store no esté vacío
    if vector_store._collection.count() == 0:
        raise ValueError("El vector store está vacío. No se cargaron documentos.")

    # Realizar una búsqueda de prueba
    try:
        query = "Area de planeamiento de ms4m"  # Cambia esto por una consulta relevante para tu caso
        resultados = vector_store.similarity_search(query, k=3)
        print(f"Búsqueda de prueba realizada. Resultados encontrados: {len(resultados)}")
    except Exception as e:
        print(f"Error al realizar la búsqueda de prueba: {e}")
        raise

    return vector_store


## PROCESO PARA CARGAR DOCUMENTOS AL VECTOR STORE
## 1.Carga solo archivos locales 
ruta_files = [
    #"pdf_solos/ASIS_REPORTE_EJECUTIVO.pdf",  # Ruta a un archivo PDF
    "ppt_solos",  # Ruta a una carpeta de PPTs
    "word_solos",  # Ruta a una carpeta de Words
    "pdf_carpetas/pdfs_varios"  # Ruta a una carpeta de PDFs
#    "json_solos"
]

## Carga Solo API de JIRA
config_jira = {
    'jql': 'project = MS4M AND resolutiondate >= "2024-01-01" AND resolutiondate <= "2025-03-31" AND status = "CLOSED" ',
    'fields': ["key", "summary", "status", "resolution", "resolutiondate", "description", "comment", "customfield_11281",
      "customfield_11228", "customfield_10975"]
}

### 2.Cargar documentos
#Opcion 1
#documentos = cargar_documentos(ruta_files)

#Opcion 2
#documentos = cargar_documentos([], jira_config=config_jira)

#Opcion 3
### Cargar documentos
#documentos = cargar_documentos(
#    rutas_archivos=ruta_files,
#    jira_config=config_jira
#)

# Crear el vector store
#vector_store = crear_vectorstore(documentos)
